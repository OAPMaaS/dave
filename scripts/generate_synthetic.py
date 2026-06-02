#!/usr/bin/env python3
"""
generate_synthetic.py — generates synthetic corporate documents for DAVE testing.

Outputs 8 documents (DOCX, PDF, PPTX) with deliberate compliance violations
plus a manifest.json with expected_findings per document.

Usage:
    python3 generate_synthetic.py
    python3 generate_synthetic.py --out-dir /custom/path
"""

import argparse
import json
import os
import shutil
import io
from pathlib import Path
from datetime import datetime

_REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = Path(os.getenv("TEST_DOCS_DIR",  str(_REPO_ROOT / "test_docs")))
WEB_DIR = Path(os.getenv("WEB_DATA_DIR",   str(_REPO_ROOT / "test_docs")))


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def s(text: str) -> str:
    """Sanitize to latin-1 safe for fpdf2 Helvetica."""
    return (text
        .replace("—", "-").replace("–", "-")
        .replace("‘", "'").replace("’", "'")
        .replace("“", '"').replace("”", '"')
        .replace("•", "*").replace("…", "...")
        .encode("latin-1", errors="replace").decode("latin-1"))


# ---------------------------------------------------------------------------
# DOCX generators
# ---------------------------------------------------------------------------

def make_contrato_pii(out: Path) -> list[str]:
    """contrato_juan_garcia_2024.docx — ES — PII expuesta."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    cp = doc.core_properties
    cp.title  = "Contrato de Trabajo Indefinido"
    cp.author = "Departamento de RRHH"
    cp.subject = "Contrato laboral empleado"
    cp.keywords = "contrato, RRHH, laboral"
    cp.created = datetime(2024, 3, 15)
    cp.modified = datetime(2024, 3, 15)

    # Cabecera con nombre completo — violación PII
    h = doc.add_heading("Juan García López — Contrato Indefinido", level=1)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("Fecha: 15 de marzo de 2024")
    doc.add_paragraph("Empresa: Omniaccess Technologies S.L.")
    doc.add_paragraph()

    doc.add_heading("1. Datos del Empleado", level=2)
    t = doc.add_table(rows=6, cols=2)
    t.style = "Table Grid"
    data = [
        ("Nombre completo",  "Juan García López"),
        ("DNI",              "12345678A"),           # PII: DNI
        ("Teléfono",         "+34 612 345 678"),     # PII: phone
        ("Email personal",   "juan.garcia@gmail.com"), # PII: email
        ("Dirección",        "Calle Mayor 42, 3ºB, 28013 Madrid"),
        ("Fecha nacimiento", "14/07/1988"),
    ]
    for i, (k, v) in enumerate(data):
        t.rows[i].cells[0].text = k
        t.rows[i].cells[1].text = v

    doc.add_paragraph()
    doc.add_heading("2. Condiciones Laborales", level=2)
    doc.add_paragraph("Puesto: Desarrollador Senior Backend")
    doc.add_paragraph("Salario bruto anual: 52.000 EUR")
    doc.add_paragraph("Jornada: 40 horas semanales")
    doc.add_paragraph("Incorporación: 1 de abril de 2024")
    doc.add_paragraph()

    doc.add_heading("3. Cláusulas", level=2)
    for cl in [
        "El empleado se compromete a mantener la confidencialidad de la información de la empresa.",
        "El período de prueba es de 6 meses según lo establecido en el Estatuto de los Trabajadores.",
        "Cualquier modificación de este contrato deberá realizarse por escrito y con acuerdo de ambas partes.",
    ]:
        doc.add_paragraph(cl, style="List Bullet")

    doc.add_paragraph()
    doc.add_heading("4. Firmas", level=2)
    t2 = doc.add_table(rows=2, cols=2)
    t2.style = "Table Grid"
    t2.rows[0].cells[0].text = "Por la empresa:"
    t2.rows[0].cells[1].text = "El empleado:"
    t2.rows[1].cells[0].text = "María Sánchez (RRHH)"
    t2.rows[1].cells[1].text = "Juan García López"

    doc.save(out)
    return ["pii_dni", "pii_phone", "pii_email", "pii_fullname_header"]


def make_gdpr_policy_draft(out: Path) -> list[str]:
    """gdpr_policy_draft.docx — EN — placeholders, missing author/date."""
    from docx import Document

    doc = Document()
    cp = doc.core_properties
    cp.title  = "GDPR Data Protection Policy"
    cp.author = ""                          # violation: empty author
    cp.subject = "Data protection policy draft"
    cp.created = datetime(2026, 5, 1)
    cp.modified = datetime(2026, 5, 1)

    doc.add_heading("GDPR Data Protection Policy", level=1)

    meta_table = doc.add_table(rows=4, cols=2)
    meta_table.style = "Table Grid"
    rows = [
        ("Version",       "0.1 DRAFT"),
        ("Author",        "[TBD]"),             # violation: placeholder
        ("Review date",   "[INSERT DATE]"),     # violation: placeholder
        ("Approved by",   "[PENDING APPROVAL]"),# violation: placeholder
    ]
    for i, (k, v) in enumerate(rows):
        meta_table.rows[i].cells[0].text = k
        meta_table.rows[i].cells[1].text = v

    doc.add_paragraph()
    doc.add_heading("1. Purpose", level=2)
    doc.add_paragraph(
        "This policy establishes the framework for data protection compliance "
        "within [COMPANY NAME] in accordance with EU Regulation 2016/679 (GDPR)."
    )

    doc.add_heading("2. Scope", level=2)
    doc.add_paragraph("This policy applies to all employees, contractors, and third parties who process personal data on behalf of [COMPANY NAME].")

    doc.add_heading("3. Data Categories", level=2)
    doc.add_paragraph("TODO: Define categories of personal data processed by the organisation.")  # placeholder
    doc.add_paragraph("TODO: Add retention periods per data category.")                            # placeholder

    doc.add_heading("4. Data Subject Rights", level=2)
    for right in [
        "Right of access (Art. 15 GDPR)",
        "Right to rectification (Art. 16 GDPR)",
        "Right to erasure (Art. 17 GDPR)",
        "Right to data portability (Art. 20 GDPR)",
    ]:
        doc.add_paragraph(right, style="List Bullet")

    doc.add_heading("5. Breach Notification", level=2)
    doc.add_paragraph("In the event of a data breach, [RESPONSIBLE PERSON - TBD] must be notified within 72 hours.")

    doc.add_heading("6. Review", level=2)
    doc.add_paragraph("This policy will be reviewed annually. Next review: [INSERT DATE].")

    doc.save(out)
    return ["placeholder_values", "missing_author", "missing_review_date"]


def make_spec_autenticacion(out: Path) -> list[str]:
    """spec_autenticacion_v0.docx — ES — versioning incompleto, placeholders."""
    from docx import Document

    doc = Document()
    cp = doc.core_properties
    cp.title   = "Especificacion de Autenticacion - Sistema DAVE"
    cp.author  = "pendiente"                   # violation: placeholder
    cp.subject = "Especificacion tecnica"
    cp.created = datetime(2026, 5, 28)
    cp.modified = datetime(2026, 5, 28)

    doc.add_heading("Especificacion Tecnica: Autenticacion", level=1)

    meta = doc.add_table(rows=5, cols=2)
    meta.style = "Table Grid"
    for i, (k, v) in enumerate([
        ("Version",     "TBD"),                # violation: placeholder
        ("Autor",       "pendiente"),           # violation: placeholder
        ("Fecha",       "TODO"),               # violation: placeholder
        ("Estado",      "BORRADOR"),
        ("Proyecto",    "DAVE - Hackathon 2026"),
    ]):
        meta.rows[i].cells[0].text = k
        meta.rows[i].cells[1].text = v

    doc.add_paragraph()
    doc.add_heading("1. Objetivo", level=2)
    doc.add_paragraph(
        "Este documento describe el sistema de autenticacion para la API de DAVE. "
        "Define los endpoints, flujos JWT y politica de sesiones."
    )

    doc.add_heading("2. Endpoints", level=2)
    t = doc.add_table(rows=4, cols=3)
    t.style = "Table Grid"
    for i, (method, path, desc) in enumerate([
        ("Metodo",  "Endpoint",          "Descripcion"),
        ("POST",    "/api/v1/auth/login",  "Genera JWT con credenciales"),
        ("POST",    "/api/v1/auth/refresh","Renueva token expirado"),
        ("DELETE",  "/api/v1/auth/logout", "Invalida sesion activa"),
    ]):
        t.rows[i].cells[0].text = method
        t.rows[i].cells[1].text = path
        t.rows[i].cells[2].text = desc

    doc.add_paragraph()
    doc.add_heading("3. Seguridad", level=2)
    doc.add_paragraph("A completar — ver con Nacho los requisitos de seguridad.")  # placeholder
    doc.add_paragraph("TODO: anadir politica de rate limiting.")                    # placeholder

    doc.add_heading("4. Dependencias", level=2)
    for dep in ["python-jose >= 3.3.0", "passlib[bcrypt] >= 1.7.4", "fastapi >= 0.110.0"]:
        doc.add_paragraph(dep, style="List Bullet")

    # Missing section: "Tabla de Cambios" (required by estructura_specs.md)

    doc.save(out)
    return ["placeholder_values", "missing_version", "missing_author", "missing_changelog"]


def make_contrato_limpio(out: Path) -> list[str]:
    """2026-06-01_contrato_limpio.docx — ES — sin violaciones."""
    from docx import Document

    doc = Document()
    cp = doc.core_properties
    cp.title   = "Contrato de Prestacion de Servicios"
    cp.author  = "Departamento Legal"
    cp.subject = "Contrato de servicios profesionales"
    cp.keywords = "contrato, servicios, legal"
    cp.created = datetime(2026, 6, 1)
    cp.modified = datetime(2026, 6, 1)

    doc.add_heading("Contrato de Prestacion de Servicios", level=1)

    meta = doc.add_table(rows=4, cols=2)
    meta.style = "Table Grid"
    for i, (k, v) in enumerate([
        ("Version",       "1.0"),
        ("Autor",         "Departamento Legal"),
        ("Fecha",         "2026-06-01"),
        ("Estado",        "Aprobado"),
    ]):
        meta.rows[i].cells[0].text = k
        meta.rows[i].cells[1].text = v

    doc.add_paragraph()
    doc.add_heading("1. Partes", level=2)
    doc.add_paragraph(
        "De una parte, EMPRESA CONTRATANTE S.L. (en adelante 'el Cliente'), "
        "con CIF B-12345678, representada por su Director General."
    )
    doc.add_paragraph(
        "De otra parte, PROVEEDOR DE SERVICIOS S.L. (en adelante 'el Proveedor'), "
        "con CIF B-87654321, representada por su Administrador Unico."
    )

    doc.add_heading("2. Objeto del Contrato", level=2)
    doc.add_paragraph(
        "El Proveedor se compromete a prestar servicios de consultoria tecnologica "
        "segun las especificaciones del Anexo I adjunto a este contrato."
    )

    doc.add_heading("3. Duracion", level=2)
    doc.add_paragraph("El contrato tiene una duracion de 12 meses desde la fecha de firma, "
                       "prorrogable por periodos anuales mediante acuerdo expreso.")

    doc.add_heading("4. Precio y Forma de Pago", level=2)
    doc.add_paragraph("El importe total es de 24.000 EUR + IVA, facturado mensualmente.")

    doc.add_heading("5. Confidencialidad", level=2)
    doc.add_paragraph("Ambas partes se comprometen a mantener la mas estricta confidencialidad "
                       "sobre la informacion intercambiada durante la vigencia del contrato.")

    doc.add_heading("6. Tabla de Cambios", level=2)
    t = doc.add_table(rows=2, cols=3)
    t.style = "Table Grid"
    for i, (ver, fecha, cambio) in enumerate([
        ("Version", "Fecha",      "Descripcion"),
        ("1.0",     "2026-06-01", "Version inicial aprobada"),
    ]):
        t.rows[i].cells[0].text = ver
        t.rows[i].cells[1].text = fecha
        t.rows[i].cells[2].text = cambio

    doc.save(out)
    return []


# ---------------------------------------------------------------------------
# PDF generators
# ---------------------------------------------------------------------------

def make_informe_ventas_pdf(out: Path) -> list[str]:
    """informe_ventas_Q1_2026.pdf — ES — naming incorrecto, falta portada."""
    from fpdf import FPDF

    class PDF(FPDF):
        def header(self):
            self.set_font("Helvetica", "B", 9)
            self.set_text_color(100, 100, 100)
            self.cell(0, 8, s("Omniaccess Technologies - Informe Interno"), align="R")
            self.ln(4)
            self.set_draw_color(200, 200, 200)
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(4)
        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(150, 150, 150)
            self.cell(0, 10, f"Pagina {self.page_no()}", align="C")

    pdf = PDF()
    pdf.set_author("Equipo de Ventas")
    pdf.set_title(s("Informe de Ventas Q1 2026"))
    pdf.set_creator("DAVE Synthetic Generator")
    pdf.add_page()

    # No portada, va directo al contenido — violacion estructura_informes.md
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, s("Informe de Ventas Q1 2026"), ln=True, align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 8, s("Periodo: Enero - Marzo 2026"), ln=True, align="C")
    pdf.ln(8)
    pdf.set_text_color(0, 0, 0)

    # Sin resumen ejecutivo — violacion estructura_informes.md
    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 10, s("1. Resultados Generales"), ln=True)
    pdf.set_font("Helvetica", "", 10)

    rows = [
        ("Region",       "Objetivo (EUR)", "Real (EUR)", "Cumplimiento"),
        ("Norte",        "450.000",        "487.320",    "108%"),
        ("Sur",          "380.000",        "342.100",    "90%"),
        ("Centro",       "520.000",        "558.900",    "107%"),
        ("Internacional","320.000",        "289.400",    "90%"),
        ("TOTAL",        "1.670.000",      "1.677.720",  "100,5%"),
    ]
    col_w = [50, 40, 40, 35]
    pdf.set_fill_color(230, 230, 230)
    for ri, row in enumerate(rows):
        fill = ri == 0 or ri == len(rows) - 1
        pdf.set_font("Helvetica", "B" if fill else "", 9)
        for ci, cell in enumerate(row):
            pdf.cell(col_w[ci], 7, s(cell), border=1, fill=fill)
        pdf.ln()
    pdf.ln(8)

    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 10, s("2. Top 5 Clientes"), ln=True)
    pdf.set_font("Helvetica", "", 10)

    clientes = [
        ("Cliente",        "Sector",    "Facturacion Q1"),
        ("TechCorp S.A.",  "Tecnologia","128.400 EUR"),
        ("Retail Plus",    "Comercio",  "97.200 EUR"),
        ("IndustrialGroup","Industria", "89.500 EUR"),
        ("FinanceHub",     "Finanzas",  "76.300 EUR"),
        ("MediaCo",        "Medios",    "64.100 EUR"),
    ]
    col_w2 = [70, 40, 55]
    for ri, row in enumerate(clientes):
        pdf.set_font("Helvetica", "B" if ri == 0 else "", 9)
        for ci, cell in enumerate(row):
            pdf.cell(col_w2[ci], 7, s(cell), border=1, fill=(ri==0))
        pdf.ln()
    pdf.ln(8)

    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 10, s("3. Conclusiones"), ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 6, s(
        "El Q1 2026 ha cerrado con un cumplimiento del 100,5% sobre objetivo. "
        "Las regiones Norte y Centro han superado sus metas, mientras Sur e Internacional "
        "presentan desviaciones que requieren atencion en Q2."
    ))

    # Sin indice, sin portada — violaciones estructura_informes.md
    pdf.output(str(out))
    return ["naming_convention", "missing_cover_page", "missing_executive_summary", "missing_index"]


def make_api_design_pdf(out: Path) -> list[str]:
    """api_design_FINAL_v2.pdf — EN — naming chaos, exposed emails."""
    from fpdf import FPDF

    class PDF(FPDF):
        def header(self):
            self.set_font("Helvetica", "B", 8)
            self.set_text_color(120, 120, 120)
            self.cell(0, 8, "DAVE API Design - INTERNAL", align="R")
            self.ln(4)
        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(150, 150, 150)
            self.cell(0, 10, f"Page {self.page_no()}", align="C")

    pdf = PDF()
    pdf.set_author("john.smith@omniaccess-internal.com")   # PII: exposed email
    pdf.set_title("DAVE API Design FINAL v2")
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 20)
    pdf.cell(0, 14, "DAVE REST API Design", ln=True, align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(130, 130, 130)
    pdf.cell(0, 7, "Version: FINAL_v2  |  Status: ???  |  Last modified: see git", ln=True, align="C")
    pdf.ln(10)
    pdf.set_text_color(0, 0, 0)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 9, "Contact & Ownership", ln=True)
    pdf.set_font("Helvetica", "", 10)
    # Exposed emails — PII violation
    contacts = [
        ("Tech Lead",    "john.smith@omniaccess-internal.com",  "+1 415 555 0192"),
        ("Backend Dev",  "anna.mueller@omniaccess-internal.com",""),
        ("DevOps",       "pedro.ruiz@omniaccess-internal.com",  "+34 600 123 456"),
    ]
    col_w = [45, 90, 50]
    pdf.set_fill_color(230, 230, 230)
    for ri, row in enumerate(contacts):
        pdf.set_font("Helvetica", "B" if ri == 0 else "", 9)
        for ci, cell in enumerate(row):
            pdf.cell(col_w[ci], 7, s(cell), border=1, fill=(ri==0))
        pdf.ln()
    pdf.ln(8)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 9, "Endpoints", ln=True)
    endpoints = [
        ("Method", "Path",                    "Auth",   "Description"),
        ("POST",   "/api/v1/documents/upload","JWT",    "Upload document for validation"),
        ("GET",    "/api/v1/runs/{id}",       "JWT",    "Get validation run status"),
        ("GET",    "/api/v1/runs/{id}/findings","JWT",  "List findings for a run"),
        ("PATCH",  "/api/v1/findings/{id}",   "JWT",    "Update finding status"),
        ("POST",   "/api/v1/runs/{id}/fix",   "JWT",    "Trigger ReAct fix"),
    ]
    col_w2 = [18, 70, 16, 85]
    pdf.set_font("Helvetica", "", 9)
    for ri, row in enumerate(endpoints):
        pdf.set_font("Helvetica", "B" if ri == 0 else "", 9)
        for ci, cell in enumerate(row):
            pdf.cell(col_w2[ci], 7, s(cell), border=1, fill=(ri==0))
        pdf.ln()
    pdf.ln(8)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 9, "Authentication", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 6, s(
        "All endpoints require a Bearer JWT token in the Authorization header. "
        "Tokens are issued by POST /api/v1/auth/login and expire after 24 hours."
    ))

    pdf.output(str(out))
    return ["naming_convention", "pii_email", "pii_phone", "missing_version_field"]


# ---------------------------------------------------------------------------
# PPTX generators
# ---------------------------------------------------------------------------

def make_kickoff_pptx(out: Path) -> list[str]:
    """kickoff_proyecto_alpha.pptx — ES — placeholders, PII en notas."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.core_properties.title  = "Kickoff Proyecto Alpha"
    prs.core_properties.author = "Maria Gonzalez"

    blank_layout = prs.slide_layouts[6]   # blank
    title_layout = prs.slide_layouts[0]   # title slide
    bullet_layout = prs.slide_layouts[1]  # title + content

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        return slide

    def add_content_slide(title, bullets, notes_text=None):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text
        return slide

    # Slide 1 — Title
    add_title_slide(
        "Proyecto Alpha — Kickoff",
        "Omniaccess Technologies | Junio 2026"
    )

    # Slide 2 — Objetivos con placeholders y PII en notas
    add_content_slide(
        "Objetivos del Proyecto",
        [
            "TBD — definir con el cliente en la proxima reunion",   # placeholder
            "Pendiente: alcance tecnico por confirmar",             # placeholder
            "Reducir tiempo de procesamiento documental en un X%",  # placeholder (X%)
            "Integracion con SharePoint corporativo",
        ],
        notes_text=(                                                # PII en notas
            "NOTAS PRIVADAS: Contacto cliente: Roberto Fernandez, "
            "+34 699 887 766, roberto.fernandez@clienteconfidencial.com. "
            "Budget aprobado: 85.000 EUR. No mencionar en la presentacion."
        )
    )

    # Slide 3 — Equipo
    add_content_slide(
        "Equipo del Proyecto",
        [
            "Project Manager: Maria Gonzalez",
            "Tech Lead: Nacho Perez",
            "Backend: Luca, Augusto",
            "Frontend: Marc",
        ]
    )

    # Slide 4 — Timeline con TBDs
    add_content_slide(
        "Timeline",
        [
            "Fase 1 (Junio): Setup e infraestructura — TODO: confirmar fechas",  # placeholder
            "Fase 2 (Julio): Desarrollo core — A definir",                       # placeholder
            "Fase 3 (Agosto): Testing y QA — Pendiente de recursos",             # placeholder
            "Go-live: TBD",                                                      # placeholder
        ]
    )

    # Slide 5 — Proximos pasos
    add_content_slide(
        "Proximos Pasos",
        [
            "Firmar NDA con el cliente (responsable: PENDIENTE)",   # placeholder
            "Configurar entornos de desarrollo",
            "Primera demo: fecha TBD",                               # placeholder
            "Kick-off tecnico interno: esta semana",
        ]
    )

    prs.save(str(out))
    return ["placeholder_values", "pii_in_notes", "pii_phone", "pii_email"]


def make_q1_deck_pptx(out: Path) -> list[str]:
    """Q1_Results_Deck.pptx — EN — naming convention, PII en slide de contacto."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.core_properties.title  = "Q1 2026 Results"
    prs.core_properties.author = "Finance Team"

    title_layout  = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_content_slide(title, bullets):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b

    add_title_slide("Q1 2026 Results", "Omniaccess Technologies | April 2026")

    add_content_slide(
        "Revenue Overview",
        [
            "Total Q1 Revenue: EUR 1,677,720",
            "vs Target: +0.5% (EUR 1,670,000)",
            "North Region: 108% of target",
            "International: 90% of target — action required",
        ]
    )

    add_content_slide(
        "Key Wins",
        [
            "TechCorp contract renewal: EUR 128,400",
            "3 new enterprise clients onboarded",
            "NPS score: 72 (up from 68 in Q4 2025)",
        ]
    )

    # Slide with exposed personal contacts — PII violation
    add_content_slide(
        "Finance Contacts",
        [
            "CFO: Sarah Johnson — sarah.johnson@omniaccess-internal.com — +1 212 555 0147",  # PII
            "Controller: Luis Mora — luis.mora@omniaccess-internal.com — +34 91 555 0234",   # PII
            "AR Team: ar-team@omniaccess-internal.com",                                       # PII
        ]
    )

    add_content_slide(
        "Q2 Outlook",
        [
            "Target: EUR 1,750,000",
            "Focus: South Region recovery plan",
            "New initiative: DAVE document compliance (internal)",
            "Headcount: TBD pending HR approval",                   # placeholder
        ]
    )

    prs.save(str(out))
    return ["naming_convention", "pii_email", "pii_phone"]


# ---------------------------------------------------------------------------
# Manifest
# ---------------------------------------------------------------------------

def build_manifest(documents: list[dict]) -> dict:
    return {
        "generated": datetime.now().strftime("%Y-%m-%d"),
        "description": "Documentos sinteticos corporativos para testing de DAVE",
        "total": len(documents),
        "with_findings": sum(1 for d in documents if d["expected_findings"]),
        "clean": sum(1 for d in documents if not d["expected_findings"]),
        "documents": documents,
    }


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out-dir", default=str(OUT_DIR))
    args = parser.parse_args()

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    generators = [
        # (filename,                         generator fn,             lang, doc_type)
        ("contrato_juan_garcia_2024.docx",   make_contrato_pii,        "es", "contrato"),
        ("gdpr_policy_draft.docx",           make_gdpr_policy_draft,   "en", "policy"),
        ("spec_autenticacion_v0.docx",       make_spec_autenticacion,  "es", "spec"),
        ("2026-06-01_contrato_limpio.docx",  make_contrato_limpio,     "es", "contrato"),
        ("informe_ventas_Q1_2026.pdf",       make_informe_ventas_pdf,  "es", "informe"),
        ("api_design_FINAL_v2.pdf",          make_api_design_pdf,      "en", "spec"),
        ("kickoff_proyecto_alpha.pptx",      make_kickoff_pptx,        "es", "presentacion"),
        ("Q1_Results_Deck.pptx",             make_q1_deck_pptx,        "en", "presentacion"),
    ]

    documents = []
    for filename, fn, lang, doc_type in generators:
        out_path = out_dir / filename
        try:
            findings = fn(out_path)
            size = out_path.stat().st_size
            print(f"  {'OK' if not findings else 'PII'} {filename} ({size//1024}KB) — {len(findings)} findings")
            documents.append({
                "filename":         filename,
                "format":           filename.rsplit(".", 1)[-1],
                "lang":             lang,
                "doc_type":         doc_type,
                "size_bytes":       size,
                "expected_findings": findings,
                "has_violations":   bool(findings),
            })
        except Exception as e:
            print(f"  ERR {filename}: {e}")
            documents.append({
                "filename": filename, "format": filename.rsplit(".", 1)[-1],
                "lang": lang, "doc_type": doc_type,
                "size_bytes": 0, "expected_findings": [], "has_violations": False,
                "error": str(e),
            })

    manifest = build_manifest(documents)
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False))
    print(f"\n  manifest.json -> {manifest_path}")

    # Copy to web directory for public access
    WEB_DIR.mkdir(parents=True, exist_ok=True)
    for f in out_dir.iterdir():
        shutil.copy2(f, WEB_DIR / f.name)
    print(f"  Copied {len(list(out_dir.iterdir()))} files to {WEB_DIR}")

    print(f"\nDone: {manifest['with_findings']} with findings, {manifest['clean']} clean")


if __name__ == "__main__":
    main()
