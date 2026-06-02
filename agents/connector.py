"""
Connector Agent — reads PDF, DOCX, XLSX, or PPTX and returns a normalized DocumentSchema.

Output schema:
{
    "filename": str,
    "format": "pdf" | "docx" | "xlsx" | "pptx",
    "text": str,                        # full plain text
    "metadata": {
        "title": str | None,
        "author": str | None,
        "created": str | None,          # ISO 8601
        "modified": str | None,         # ISO 8601
        "pages": int | None,
        "slide_count": int | None,      # PPTX only
        "sheet_names": list[str] | None # XLSX only
    },
    "sections": [{"heading": str, "content": str}],
    "tables": [list[list[str]]],        # each table is a list of rows
    "images_count": int,
    "notes": str | None,                # speaker notes (PPTX only)
    "error": str | None                 # set if parsing failed partially
}
"""

import io
import os
from datetime import datetime
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def load_document(source: str | bytes | Path, filename: str | None = None) -> dict:
    """
    Parse a document and return a normalized DocumentSchema dict.

    source: file path (str/Path) or raw bytes
    filename: required when source is bytes — used to detect format
    """
    if isinstance(source, (str, Path)):
        path = Path(source)
        filename = filename or path.name
        data = path.read_bytes()
    else:
        data = source
        if not filename:
            raise ValueError("filename required when source is bytes")

    ext = Path(filename).suffix.lower().lstrip(".")
    fmt_map = {
        "pdf": "pdf", "docx": "docx", "doc": "docx",
        "xlsx": "xlsx", "xls": "xlsx",
        "pptx": "pptx", "ppt": "pptx",
    }
    fmt = fmt_map.get(ext)

    if fmt == "pdf":
        return _parse_pdf(data, filename)
    elif fmt == "docx":
        return _parse_docx(data, filename)
    elif fmt == "xlsx":
        return _parse_xlsx(data, filename)
    elif fmt == "pptx":
        return _parse_pptx(data, filename)
    else:
        return _empty_schema(filename, "unknown", f"Unsupported format: .{ext}")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _parse_pdf(data: bytes, filename: str) -> dict:
    import pdfplumber

    schema = _empty_schema(filename, "pdf")
    pages_text = []
    tables = []

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            schema["metadata"]["pages"] = len(pdf.pages)

            meta = pdf.metadata or {}
            schema["metadata"]["title"] = meta.get("Title") or None
            schema["metadata"]["author"] = meta.get("Author") or None
            schema["metadata"]["created"] = _pdf_date(meta.get("CreationDate"))
            schema["metadata"]["modified"] = _pdf_date(meta.get("ModDate"))

            for page in pdf.pages:
                text = page.extract_text() or ""
                pages_text.append(text)

                for table in page.extract_tables() or []:
                    clean = [[str(c or "").strip() for c in row] for row in table if any(c for c in row)]
                    if clean:
                        tables.append(clean)

        schema["text"] = "\n".join(pages_text).strip()
        schema["tables"] = tables
        schema["sections"] = _sections_from_text(schema["text"])
        schema["images_count"] = _count_pdf_images(data)

    except Exception as e:
        schema["error"] = f"PDF parse error: {e}"

    return schema


def _pdf_date(raw: str | None) -> str | None:
    if not raw:
        return None
    raw = str(raw).strip()
    # PDF date format: D:YYYYMMDDHHmmSSOHH'mm'
    if raw.startswith("D:"):
        raw = raw[2:]
    try:
        dt = datetime.strptime(raw[:14], "%Y%m%d%H%M%S")
        return dt.isoformat()
    except Exception:
        return raw[:14] if len(raw) >= 8 else None


def _count_pdf_images(data: bytes) -> int:
    try:
        import pdfplumber
        count = 0
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                count += len(page.images or [])
        return count
    except Exception:
        return 0


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _parse_docx(data: bytes, filename: str) -> dict:
    from docx import Document
    from docx.oxml.ns import qn

    schema = _empty_schema(filename, "docx")

    try:
        doc = Document(io.BytesIO(data))

        # Metadata
        cp = doc.core_properties
        schema["metadata"]["title"] = cp.title or None
        schema["metadata"]["author"] = cp.author or None
        schema["metadata"]["created"] = cp.created.isoformat() if cp.created else None
        schema["metadata"]["modified"] = cp.modified.isoformat() if cp.modified else None

        # Sections + full text
        sections = []
        current_heading = None
        current_lines = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading"):
                if current_heading is not None or current_lines:
                    sections.append({
                        "heading": current_heading or "",
                        "content": "\n".join(current_lines).strip(),
                    })
                current_heading = text
                current_lines = []
            else:
                current_lines.append(text)

        if current_heading is not None or current_lines:
            sections.append({
                "heading": current_heading or "",
                "content": "\n".join(current_lines).strip(),
            })

        schema["sections"] = sections

        # Tables — extract and include cell text in full text
        table_text_parts = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
                table_text_parts.append(" | ".join(c for c in cells if c))
            if rows:
                schema["tables"].append(rows)

        section_text = "\n".join(
            (f"# {s['heading']}\n{s['content']}" if s["heading"] else s["content"])
            for s in sections
        )
        schema["text"] = (section_text + "\n" + "\n".join(table_text_parts)).strip()

        # Images
        schema["images_count"] = sum(
            1 for rel in doc.part.rels.values()
            if "image" in rel.reltype
        )

    except Exception as e:
        schema["error"] = f"DOCX parse error: {e}"

    return schema


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def _parse_xlsx(data: bytes, filename: str) -> dict:
    import openpyxl

    schema = _empty_schema(filename, "xlsx")

    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)

        schema["metadata"]["sheet_names"] = wb.sheetnames

        props = wb.properties
        schema["metadata"]["title"] = props.title or None
        schema["metadata"]["author"] = (props.creator or props.lastModifiedBy) or None
        schema["metadata"]["created"] = props.created.isoformat() if props.created else None
        schema["metadata"]["modified"] = props.modified.isoformat() if props.modified else None

        all_text = []
        tables = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            sheet_text_lines = [f"[Sheet: {sheet_name}]"]

            for row in ws.iter_rows(values_only=True):
                if all(v is None for v in row):
                    continue
                cells = [str(v).strip() if v is not None else "" for v in row]
                rows.append(cells)
                sheet_text_lines.append("\t".join(cells))

            if rows:
                tables.append(rows)
                all_text.extend(sheet_text_lines)

            # Sections = named ranges or sheet name as heading
            schema["sections"].append({
                "heading": sheet_name,
                "content": "\n".join(sheet_text_lines[1:]),
            })

        schema["text"] = "\n".join(all_text).strip()
        schema["tables"] = tables

    except Exception as e:
        schema["error"] = f"XLSX parse error: {e}"

    return schema


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _parse_pptx(data: bytes, filename: str) -> dict:
    from pptx import Presentation as PRS
    from pptx.util import Pt

    schema = _empty_schema(filename, "pptx")

    try:
        prs = PRS(io.BytesIO(data))

        cp = prs.core_properties
        schema["metadata"]["title"]       = cp.title or None
        schema["metadata"]["author"]      = cp.author or None
        schema["metadata"]["created"]     = cp.created.isoformat() if cp.created else None
        schema["metadata"]["modified"]    = cp.modified.isoformat() if cp.modified else None
        schema["metadata"]["slide_count"] = len(prs.slides)

        sections = []
        all_text = []
        all_notes = []
        images = 0

        for i, slide in enumerate(prs.slides, 1):
            slide_title = ""
            slide_lines = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        if shape.shape_type == 13:  # placeholder title
                            slide_title = text
                        else:
                            slide_lines.append(text)

                if hasattr(shape, "image"):
                    images += 1

                if shape.shape_type == 19:  # table
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        schema["tables"].append(rows)

            # Try to get title from title placeholder
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()
                slide_lines = [l for l in slide_lines if l != slide_title]

            heading = slide_title or f"Slide {i}"
            content = "\n".join(slide_lines)
            sections.append({"heading": heading, "content": content})
            all_text.append(f"[{heading}]\n{content}" if content else f"[{heading}]")

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    all_notes.append(f"[{heading}] {notes}")

        schema["sections"]     = sections
        schema["text"]         = "\n\n".join(all_text).strip()
        schema["notes"]        = "\n".join(all_notes) if all_notes else None
        schema["images_count"] = images

    except Exception as e:
        schema["error"] = f"PPTX parse error: {e}"

    return schema


def _empty_schema(filename: str, fmt: str, error: str | None = None) -> dict:
    return {
        "filename": filename,
        "format": fmt,
        "text": "",
        "metadata": {
            "title": None,
            "author": None,
            "created": None,
            "modified": None,
            "pages": None,
            "slide_count": None,
            "sheet_names": None,
        },
        "sections": [],
        "tables": [],
        "images_count": 0,
        "notes": None,
        "error": error,
    }


def _sections_from_text(text: str) -> list[dict]:
    """
    Heuristic section splitter for PDF text: lines in ALL CAPS or ending with ':'
    that look like headings, followed by their body.
    """
    lines = text.split("\n")
    sections = []
    current_heading = None
    current_lines = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        is_heading = (
            len(stripped) < 80
            and (stripped.isupper() or stripped.endswith(":"))
            and not stripped.endswith(".")
        )
        if is_heading:
            if current_heading is not None or current_lines:
                sections.append({
                    "heading": current_heading or "",
                    "content": "\n".join(current_lines).strip(),
                })
            current_heading = stripped.rstrip(":")
            current_lines = []
        else:
            current_lines.append(stripped)

    if current_heading is not None or current_lines:
        sections.append({
            "heading": current_heading or "",
            "content": "\n".join(current_lines).strip(),
        })

    return sections


# ---------------------------------------------------------------------------
# CLI — python3 connector.py <file>
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys
    import json

    if len(sys.argv) < 2:
        print("Usage: python3 connector.py <file.pdf|.docx|.xlsx>")
        sys.exit(1)

    result = load_document(sys.argv[1])
    result_preview = {**result, "text": result["text"][:500] + ("..." if len(result["text"]) > 500 else "")}
    print(json.dumps(result_preview, indent=2, ensure_ascii=False, default=str))
