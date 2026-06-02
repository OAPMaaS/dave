"""
Generate a realistic demo corpus for the AI-Readiness Auditor.

Target: 45 files, ~17 flagged (~38%) — a believable "this company has a data
problem" story. Majority are real Office files (.docx/.xlsx/.pptx); 4 deliberate
orphan formats (.txt/.json/.csv) for the format-signal demo.

Output: domain/demo_corpus/files/  (never scanned alongside this script)

Run:
    python -m domain.demo_corpus.generate_sample
    python -m domain.demo_corpus.generate_sample --output /tmp/corp_docs
"""
from __future__ import annotations

import argparse
import csv as csv_mod
import json
import os
from datetime import datetime, timedelta, timezone
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt as PPTXPt

# ── Low-level helpers ─────────────────────────────────────────────────────────

def _ts(days_ago: int) -> float:
    return (datetime.now(tz=timezone.utc) - timedelta(days=days_ago)).timestamp()


def _backdate(path: str | Path, days: int) -> None:
    ts = _ts(days)
    os.utime(str(path), (ts, ts))


def _gov_table(doc, owner: str, classification: str,
               effective_date: str = "", review_date: str = "") -> None:
    """2-column governance header table; text format matches governance.py regexes."""
    rows = [("Owner:", owner), ("Classification:", classification)]
    if effective_date:
        rows.append(("Effective Date:", effective_date))
    if review_date:
        rows.append(("Review Date:", review_date))
    tbl = doc.add_table(rows=len(rows), cols=2)
    tbl.style = "Table Grid"
    for i, (label, val) in enumerate(rows):
        tbl.cell(i, 0).text = label
        tbl.cell(i, 1).text = val
    doc.add_paragraph("")


def _h(doc, text: str, level: int = 1) -> None:
    doc.add_heading(text, level=level)


def _p(doc, text: str) -> None:
    doc.add_paragraph(text)


# ── DOCX builders ─────────────────────────────────────────────────────────────

def _policy(path: Path, title: str, owner: str, classification: str,
            eff: str, rev: str, purpose: str, scope: str,
            extra_body: str = "") -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    _gov_table(doc, owner, classification, eff, rev)
    _h(doc, "Purpose"); _p(doc, purpose)
    _h(doc, "Scope"); _p(doc, scope)
    _h(doc, "Owner"); _p(doc, f"Policy owner: {owner}.")
    _h(doc, "Effective Date"); _p(doc, eff)
    _h(doc, "Review Date"); _p(doc, rev)
    if extra_body:
        _p(doc, extra_body)
    doc.save(path)


def _procedure(path: Path, title: str, owner: str, classification: str,
               last_updated: str, purpose: str, scope: str, steps: list[str]) -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    _gov_table(doc, owner, classification, review_date=last_updated)
    _h(doc, "Purpose"); _p(doc, purpose)
    _h(doc, "Scope"); _p(doc, scope)
    _h(doc, "Owner"); _p(doc, owner)
    _h(doc, "Steps")
    for i, step in enumerate(steps, 1):
        _p(doc, f"{i}. {step}")
    _h(doc, "Last Updated"); _p(doc, last_updated)
    doc.save(path)


def _contract(path: Path, title: str, owner: str, parties: str,
              eff: str, term: str, termination: str, counterparty: str,
              expiry: str = "", classification: str = "confidential") -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    rows = [("Owner:", owner), ("Classification:", classification),
            ("Counterparty:", counterparty), ("Effective Date:", eff)]
    if expiry:
        rows.append(("Expiry Date:", expiry))
    tbl = doc.add_table(rows=len(rows), cols=2)
    tbl.style = "Table Grid"
    for i, (lbl, val) in enumerate(rows):
        tbl.cell(i, 0).text = lbl
        tbl.cell(i, 1).text = val
    doc.add_paragraph("")
    _h(doc, "Parties"); _p(doc, parties)
    _h(doc, "Term"); _p(doc, term)
    _h(doc, "Effective Date"); _p(doc, eff)
    _h(doc, "Termination"); _p(doc, termination)
    _h(doc, "Signatures"); _p(doc, "Authorised signatories as listed above.")
    doc.save(path)


def _charter(path: Path, title: str, owner: str, classification: str,
             objective: str, scope: str, sponsor: str,
             milestones: list[str], budget: str) -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    _gov_table(doc, owner, classification)
    _h(doc, "Objective"); _p(doc, objective)
    _h(doc, "Scope"); _p(doc, scope)
    _h(doc, "Sponsor"); _p(doc, sponsor)
    _h(doc, "Milestones")
    for m in milestones:
        _p(doc, f"• {m}")
    _h(doc, "Budget"); _p(doc, budget)
    doc.save(path)


def _status(path: Path, title: str, owner: str, classification: str,
            period: str, status: str, risks: str, next_steps: str) -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    _gov_table(doc, owner, classification, review_date=period)
    _h(doc, "Period"); _p(doc, period)
    _h(doc, "Status"); _p(doc, status)
    _h(doc, "Risks"); _p(doc, risks)
    _h(doc, "Next Steps"); _p(doc, next_steps)
    doc.save(path)


def _okr(path: Path, title: str, owner: str, classification: str,
         objective: str, key_results: list[str], target: str) -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    _gov_table(doc, owner, classification)
    _h(doc, "Objective"); _p(doc, objective)
    _h(doc, "Key Results")
    for kr in key_results:
        _p(doc, f"• {kr}")
    _h(doc, "Owner"); _p(doc, owner)
    _h(doc, "Target"); _p(doc, target)
    doc.save(path)


def _raid(path: Path, title: str, owner: str, classification: str,
          risks: str, assumptions: str, issues: str, deps: str) -> None:
    doc = Document()
    doc.core_properties.author = owner
    doc.core_properties.title = title
    _h(doc, title, 0)
    _gov_table(doc, owner, classification)
    _h(doc, "Risks"); _p(doc, risks)
    _h(doc, "Assumptions"); _p(doc, assumptions)
    _h(doc, "Issues"); _p(doc, issues)
    _h(doc, "Dependencies"); _p(doc, deps)
    doc.save(path)


# ── XLSX builders ─────────────────────────────────────────────────────────────

def _data_dict(path: Path, title: str, fields: list[dict],
               owner: str = "", steward: str = "",
               classification: str = "", retention: str = "") -> None:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data Dictionary"
    if owner:
        wb.properties.creator = owner
    wb.properties.title = title
    row = 1
    # Governance block — parsed by governance.py regex
    if owner:
        ws.cell(row, 1, "Owner:"); ws.cell(row, 2, owner); row += 1
    if classification:
        ws.cell(row, 1, "Classification:"); ws.cell(row, 2, classification); row += 1
    if retention:
        ws.cell(row, 1, "Retention:"); ws.cell(row, 2, retention); row += 1
    if steward:
        ws.cell(row, 1, "Steward:"); ws.cell(row, 2, steward); row += 1
    if any([owner, classification, retention, steward]):
        row += 1  # blank spacer
    # Column headers — checked by standards.py for "Field Name", "Data Type", etc.
    headers = ["Field Name", "Data Type", "Description", "Owner", "Classification",
               "PII", "Example"]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row, col, h)
        cell.font = Font(bold=True)
    row += 1
    for f in fields:
        ws.cell(row, 1, f.get("name", ""))
        ws.cell(row, 2, f.get("type", ""))
        ws.cell(row, 3, f.get("desc", ""))
        ws.cell(row, 4, f.get("owner", owner))
        ws.cell(row, 5, f.get("classification", classification))
        ws.cell(row, 6, f.get("pii", "No"))
        ws.cell(row, 7, f.get("example", ""))
        row += 1
    wb.save(path)


def _data_contract(path: Path, title: str, owner: str, sla: str,
                   classification: str, retention: str,
                   schema_rows: list[dict]) -> None:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data Contract"
    wb.properties.creator = owner
    wb.properties.title = title
    r = 1
    for label, val in [("Owner:", owner), ("Classification:", classification),
                        ("SLA:", sla), ("Retention:", retention)]:
        ws.cell(r, 1, label); ws.cell(r, 2, val); r += 1
    r += 1
    # Section headers checked by standards.py
    for col, h in enumerate(["Schema", "Field", "Type", "Description"], 1):
        ws.cell(r, col, h).font = Font(bold=True)
    r += 1
    for row_data in schema_rows:
        ws.cell(r, 1, row_data.get("schema", ""))
        ws.cell(r, 2, row_data.get("field", ""))
        ws.cell(r, 3, row_data.get("type", ""))
        ws.cell(r, 4, row_data.get("desc", ""))
        r += 1
    wb.save(path)


# ── PPTX builders ─────────────────────────────────────────────────────────────

def _deck(path: Path, title: str, owner: str, classification: str,
          slides_data: list[dict]) -> None:
    prs = Presentation()
    prs.core_properties.author = owner
    prs.core_properties.title = title
    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = (
            f"Owner: {owner}\nClassification: {classification}"
        )
    # Content slides
    layout = prs.slide_layouts[1]
    for s in slides_data:
        sl = prs.slides.add_slide(layout)
        sl.shapes.title.text = s.get("title", "")
        tf = sl.placeholders[1].text_frame
        lines = s.get("bullets", [])
        if lines:
            tf.text = lines[0]
            for line in lines[1:]:
                tf.add_paragraph().text = line
    prs.save(path)


# ── Plain-text / CSV / JSON helpers ───────────────────────────────────────────

def _txt(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def _csv_file(path: Path, headers: list[str], rows: list[list]) -> None:
    with open(path, "w", newline="", encoding="utf-8") as f:
        w = csv_mod.writer(f)
        w.writerow(headers)
        w.writerows(rows)


def _json_file(path: Path, data) -> None:
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")


# ── Corpus generator ──────────────────────────────────────────────────────────

def generate(output_dir: str = "domain/demo_corpus/files") -> None:
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    count = 0

    def w(rel):
        nonlocal count
        count += 1
        return out / rel

    # ────────────────────────────────────────────────────────────────────────
    # CLEAN FILES (28) — proper format, good governance, within staleness cadence
    # ────────────────────────────────────────────────────────────────────────

    # Policies (5)
    _policy(w("data_classification_policy.docx"),
            "Data Classification Policy",
            "Jane Smith, Information Security", "confidential",
            "2024-01-15", "2026-01-15",
            "Establishes requirements for classifying and handling company data.",
            "All employees, contractors, and third parties accessing company data.")

    _policy(w("privacy_policy.docx"),
            "Privacy Policy",
            "Legal & Compliance Team", "confidential",
            "2024-03-01", "2026-03-01",
            "Describes how the company collects, uses, and protects personal data per GDPR.",
            "All data subjects whose personal data the company processes.")

    _policy(w("third_party_risk_policy.docx"),
            "Third-Party Risk Management Policy",
            "Head of Risk, Sarah Okafor", "internal",
            "2024-06-01", "2025-12-01",
            "Defines requirements for assessing and managing risks from third-party suppliers.",
            "All business units that engage third-party vendors or contractors.")

    _policy(w("business_continuity_plan.docx"),
            "Business Continuity Plan",
            "COO Office — David Chen", "confidential",
            "2025-01-10", "2026-01-10",
            "Ensures continued operation of critical business functions during disruptions.",
            "All critical business processes and IT systems.")

    _policy(w("password_policy.docx"),
            "Password and Authentication Policy",
            "IT Security — Lena Fischer", "internal",
            "2024-11-01", "2025-11-01",
            "Mandates strong authentication controls for all company accounts.",
            "All staff, contractors, and system accounts.")

    _policy(w("remote_work_policy.docx"),
            "Remote Work and Bring-Your-Own-Device Policy",
            "HR & IT Security", "internal",
            "2025-02-01", "2026-02-01",
            "Sets security expectations for staff working remotely or using personal devices.",
            "All employees and contractors working outside company premises.")

    # Procedures (5)
    _procedure(w("incident_response_procedure.docx"),
               "Incident Response Procedure",
               "CISO Office", "confidential", "2025-03-15",
               "Guides the team through detecting, containing, and recovering from security incidents.",
               "All IT staff and the security operations team.",
               ["Detect the incident via SIEM alerts or user reports.",
                "Classify severity (P1–P4) using the incident matrix.",
                "Contain the threat: isolate affected systems.",
                "Notify stakeholders per the communication tree.",
                "Eradicate root cause and restore services.",
                "Write a post-incident report within 72 hours."])

    _procedure(w("change_management_procedure.docx"),
               "IT Change Management Procedure",
               "IT Operations — Marco Rossi", "internal", "2025-01-20",
               "Controls how changes to production systems are planned, approved, and deployed.",
               "All production IT systems and the teams that maintain them.",
               ["Raise a change request in the ITSM tool.",
                "CAB reviews changes every Tuesday.",
                "Approved changes are scheduled in the maintenance window.",
                "Post-deployment verification within 30 minutes.",
                "Rollback plan must be documented before approval."])

    _procedure(w("data_retention_procedure.docx"),
               "Data Retention and Disposal Procedure",
               "Data Governance Team", "internal", "2025-04-01",
               "Defines how data is retained, archived, and securely disposed of.",
               "All business units holding personal or sensitive data.",
               ["Classify data and determine retention period from the retention schedule.",
                "Archive data to cold storage at retention threshold.",
                "Initiate secure deletion review at end-of-life.",
                "Record disposal in the data asset register.",
                "Update the retention schedule quarterly."])

    _procedure(w("backup_recovery_procedure.docx"),
               "Backup and Recovery Procedure",
               "IT Infrastructure Team", "internal", "2025-02-28",
               "Ensures data can be recovered within defined RPO/RTO targets.",
               "All production databases, file shares, and SaaS platforms.",
               ["Full backups run nightly at 02:00 UTC.",
                "Incremental backups run every 4 hours.",
                "Recovery test performed monthly — results logged in the DR register.",
                "Offsite replica updated every 24 hours.",
                "Critical systems: RPO ≤ 4 h, RTO ≤ 8 h."])

    _procedure(w("software_development_lifecycle.docx"),
               "Secure Software Development Lifecycle Procedure",
               "Engineering Lead — Priya Nair", "internal", "2025-05-01",
               "Embeds security into every phase of the software development process.",
               "All software development and DevOps teams.",
               ["Threat model new features at design stage.",
                "Static analysis (SAST) gates in CI pipeline.",
                "Dependency vulnerability scan on every build.",
                "Mandatory security review before production release.",
                "Penetration test annually or after major architecture changes."])

    # Contracts (2)
    _contract(w("supplier_agreement_acme_2024.docx"),
              "Master Services Agreement — Acme Analytics Ltd",
              "Head of Procurement", "Acme Analytics Ltd",
              "2024-09-01", "3 years from Effective Date",
              "Either party may terminate with 90 days written notice.", "Acme Analytics Ltd",
              expiry="2027-09-01")

    _contract(w("saas_master_agreement_2024.docx"),
              "SaaS Master Subscription Agreement — CloudBI Inc",
              "Head of Procurement", "CloudBI Inc",
              "2025-01-15", "12 months auto-renewing",
              "30 days written notice to terminate at renewal.", "CloudBI Inc",
              expiry="2026-01-15")

    # Project charters (2)
    _charter(w("project_charter_data_platform.docx"),
             "Project Charter — Enterprise Data Platform",
             "CDO — Fatima Al-Rashid", "internal",
             "Build a cloud-native data platform replacing the legacy Hadoop cluster.",
             "Finance, Operations, and Sales analytics workloads. Phase 1 only.",
             "Chief Data Officer, sponsored by CFO",
             ["2025-Q2: Cloud infrastructure provisioned",
              "2025-Q3: Finance data migrated",
              "2025-Q4: Operations data migrated",
              "2026-Q1: Self-serve BI launched"],
             "€1.2M capex; €180k annual opex post-go-live.")

    _charter(w("project_charter_crm_refresh.docx"),
             "Project Charter — CRM Platform Refresh",
             "Head of Sales Operations — Tom Bergmann", "internal",
             "Replace the legacy CRM with Salesforce Sales Cloud to improve pipeline visibility.",
             "Sales, Marketing, and Customer Success teams.",
             "Chief Revenue Officer",
             ["2025-Q3: Vendor selected and contracts signed",
              "2025-Q4: Data migration and parallel run",
              "2026-Q1: Go-live and legacy decommission"],
             "€450k implementation; €95k/year licence.")

    # Status reports (2) — backdated to stay inside the 30-day threshold
    _status(w("status_report_may_2025.docx"),
            "Weekly Status Report — 19 May 2025",
            "PMO — Ana Torres", "internal",
            "12 May – 19 May 2025", "GREEN — on track",
            "No critical risks. Minor delay in vendor onboarding (resolved).",
            "Complete vendor onboarding by 26 May. Kick off Phase 2 planning.")
    _backdate(out / "status_report_may_2025.docx", 18)

    _status(w("status_report_jun_2025.docx"),
            "Weekly Status Report — 2 June 2025",
            "PMO — Ana Torres", "internal",
            "26 May – 2 June 2025", "AMBER — minor delay",
            "Infrastructure provisioning 1 week behind schedule.",
            "Escalate resource request to CTO by 5 June.")
    _backdate(out / "status_report_jun_2025.docx", 5)

    # OKRs (1) — within 90-day threshold
    _okr(w("okr_q2_2025.docx"),
         "OKRs Q2 2025 — Data & Analytics",
         "CDO — Fatima Al-Rashid", "internal",
         "Become a data-driven organisation with trusted, accessible analytics.",
         ["Reduce report generation time from 48 h to 4 h — target Q2 end.",
          "Achieve 90% data quality score on CRM pipeline data.",
          "Onboard 3 new self-serve BI teams by June 30.",
          "Zero P1 data incidents in Q2."],
         "Q2 2025 (April – June)")
    _backdate(out / "okr_q2_2025.docx", 10)

    # RAID log (1)
    _raid(w("raid_log_data_platform.docx"),
          "RAID Log — Enterprise Data Platform",
          "PMO — Ana Torres", "internal",
          "Vendor lock-in risk if single cloud provider chosen. Mitigation: multi-cloud design.",
          "Cloud team has Terraform skills. CDO approved budget uplift.",
          "Legacy ETL jobs failing intermittently — investigation in progress (ticket #4421).",
          "Finance sign-off on data model required before Sprint 4.")

    # Data dictionaries (.xlsx, clean, 2)
    _data_dict(w("crm_data_dictionary.xlsx"),
               "CRM Data Dictionary",
               fields=[
                   {"name": "contact_id",    "type": "UUID",        "desc": "Primary key for the CRM contact entity.", "pii": "No"},
                   {"name": "email",          "type": "VARCHAR(255)", "desc": "Contact email address.",                  "pii": "Yes"},
                   {"name": "phone",          "type": "VARCHAR(20)",  "desc": "Primary phone number.",                   "pii": "Yes"},
                   {"name": "account_id",     "type": "UUID",        "desc": "FK to the parent account.",               "pii": "No"},
                   {"name": "lead_score",     "type": "INTEGER",     "desc": "Proprietary lead quality score 0–100.",    "pii": "No"},
                   {"name": "owner_id",       "type": "UUID",        "desc": "Assigned sales rep.",                     "pii": "No"},
                   {"name": "created_at",     "type": "TIMESTAMP",   "desc": "Record creation timestamp (UTC).",        "pii": "No"},
               ],
               owner="Data Engineering — Alice Chen",
               steward="Sales Ops — Tom Bergmann",
               classification="confidential",
               retention="7 years per GDPR Art. 17")

    _data_dict(w("finance_data_dictionary.xlsx"),
               "Finance Master Data Dictionary",
               fields=[
                   {"name": "gl_account",     "type": "VARCHAR(10)", "desc": "General ledger account code.",    "pii": "No"},
                   {"name": "cost_centre",    "type": "VARCHAR(10)", "desc": "Responsibility centre code.",     "pii": "No"},
                   {"name": "transaction_id", "type": "BIGINT",      "desc": "Unique transaction identifier.",  "pii": "No"},
                   {"name": "amount_eur",     "type": "DECIMAL(15,2)","desc": "Transaction amount in EUR.",     "pii": "No"},
                   {"name": "posting_date",   "type": "DATE",        "desc": "Journal posting date.",           "pii": "No"},
                   {"name": "vendor_id",      "type": "VARCHAR(20)", "desc": "FK to vendor master.",            "pii": "No"},
               ],
               owner="Finance Controller — Peter Müller",
               steward="Finance Controller — Peter Müller",
               classification="restricted",
               retention="10 years per tax law")

    # Data contracts (.xlsx, clean, 2)
    _data_contract(w("sales_analytics_data_contract.xlsx"),
                   "Sales Analytics Data Contract",
                   owner="Data Engineering", sla="99.5% availability; refresh ≤ 4 h",
                   classification="internal", retention="3 years",
                   schema_rows=[
                       {"schema": "sales", "field": "opportunity_id",  "type": "UUID",    "desc": "Unique opportunity identifier"},
                       {"schema": "sales", "field": "stage",           "type": "VARCHAR", "desc": "Pipeline stage"},
                       {"schema": "sales", "field": "amount",          "type": "DECIMAL", "desc": "Forecast deal value EUR"},
                       {"schema": "sales", "field": "close_date",      "type": "DATE",    "desc": "Expected close date"},
                   ])

    _data_contract(w("hr_data_contract.xlsx"),
                   "HR Data Contract — People Analytics",
                   owner="HR Data Team — Nina Kowalski", sla="99% availability; daily refresh",
                   classification="restricted", retention="7 years",
                   schema_rows=[
                       {"schema": "hr", "field": "employee_id",   "type": "UUID",    "desc": "Anonymised employee ID"},
                       {"schema": "hr", "field": "department",    "type": "VARCHAR", "desc": "Department name"},
                       {"schema": "hr", "field": "hire_date",     "type": "DATE",    "desc": "Employment start date"},
                       {"schema": "hr", "field": "tenure_months", "type": "INTEGER", "desc": "Months of service"},
                   ])

    # Decks (.pptx, clean, 2)
    _deck(w("data_strategy_q1_2025.pptx"),
          "Data & Analytics Strategy — Q1 2025", "CDO — Fatima Al-Rashid", "internal",
          [{"title": "Vision", "bullets": ["Single source of truth by end of 2025",
                                           "Self-serve analytics for all business units",
                                           "AI-ready data by 2026"]},
           {"title": "Initiatives", "bullets": ["Enterprise Data Platform (Phase 1 live Q3)",
                                                 "CRM Data Contract programme",
                                                 "Data Governance Council launch"]},
           {"title": "Metrics", "bullets": ["Data quality score: 72% → 90%",
                                             "Report latency: 48 h → 4 h",
                                             "Self-serve adoption: 18% → 60%"]}])

    _deck(w("board_fy2025_data_overview.pptx"),
          "Board Pack — FY2025 Data & Technology Overview",
          "CTO — James Obi", "restricted",
          [{"title": "Technology Investment", "bullets": ["Cloud migration: €3.2M",
                                                           "Data Platform: €1.2M",
                                                           "Cyber resilience: €800k"]},
           {"title": "Risk Summary", "bullets": ["Third-party cyber risk: AMBER",
                                                  "Data quality risk: AMBER (improving)",
                                                  "Regulatory compliance: GREEN"]}])

    print(f"  [CLEAN] {count} files so far")
    clean_count = count

    # ────────────────────────────────────────────────────────────────────────
    # FLAGGED — STALE WITH BAD CONTENT (10)
    # ────────────────────────────────────────────────────────────────────────

    # JUICY #1: Security policy — ISO 27001:2013 + PCI DSS 3.2.1, no owner
    doc = Document()
    doc.add_heading("Information Security Policy", 0)
    doc.add_paragraph(
        "Review Date: 2022-06-30\n"
        "Effective Date: 2019-09-01"
    )
    _h(doc, "Purpose")
    doc.add_paragraph(
        "This policy defines the information security requirements for all company "
        "systems and staff. Controls are aligned with ISO/IEC 27001:2013 Annex A "
        "and PCI DSS 3.2.1 requirements. Network segmentation follows Windows Server "
        "2012 Active Directory domain guidelines."
    )
    _h(doc, "Scope")
    doc.add_paragraph("All company IT systems, employees, and third parties.")
    _h(doc, "Effective Date"); doc.add_paragraph("1 September 2019")
    _h(doc, "Review Date"); doc.add_paragraph("30 June 2022")
    # Deliberately omit: Owner heading, Classification, governance table
    doc.save(out / "security_policy_2021.docx")
    _backdate(out / "security_policy_2021.docx", 1400)

    # JUICY #2: Data governance framework — DRAFT, all TBD, backdated
    doc = Document()
    doc.add_heading("Data Governance Framework", 0)
    doc.add_paragraph("DRAFT — DO NOT DISTRIBUTE")
    _h(doc, "Purpose")
    doc.add_paragraph("TBD — framework purpose to be defined after stakeholder workshops.")
    _h(doc, "Scope")
    doc.add_paragraph("[INSERT scope here — awaiting input from data owners]")
    _h(doc, "Roles and Responsibilities")
    doc.add_paragraph(
        "Data Owner: [INSERT name]\nData Steward: TODO\nGovernance Committee: TBD"
    )
    _h(doc, "Standards")
    doc.add_paragraph(
        "Compliance with ISO/IEC 27001:2013 and Data Protection Directive 95/46/EC "
        "will be assessed in Phase 2. Lorem ipsum placeholder text for now."
    )
    doc.save(out / "data_governance_framework_draft_2020.docx")
    _backdate(out / "data_governance_framework_draft_2020.docx", 1600)

    # JUICY #3 (xlsx): customer data dictionary — no owner, no steward, no classification
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Customer Data Dictionary"
    wb.properties.title = "Customer Data Dictionary v1"
    # No creator/author set (no owner in embedded metadata)
    # Column headers only — no governance block, no "Owner" or "Classification" rows
    for col, h in enumerate(["Field Name", "Data Type", "Description", "Notes"], 1):
        ws.cell(1, col, h).font = Font(bold=True)
    rows_data = [
        ("customer_id",   "UUID",        "Primary key.",               ""),
        ("email",         "VARCHAR(255)","Customer email.",             "PII — handle with care"),
        ("credit_score",  "INTEGER",     "Proprietary score 300-850.",  "Source: internal model v3"),
        ("date_of_birth", "DATE",        "Date of birth.",              "PII — GDPR Art. 9"),
        ("national_id",   "VARCHAR(20)", "Government-issued ID number.","PII — restricted"),
        ("home_address",  "TEXT",        "Full home address.",          "PII"),
        ("income_band",   "VARCHAR(20)", "Estimated income bracket.",   "Derived field"),
    ]
    for i, row in enumerate(rows_data, 2):
        for j, val in enumerate(row, 1):
            ws.cell(i, j, val)
    wb.save(out / "customer_data_dictionary_v1.xlsx")

    # JUICY #4: GDPR compliance policy — transition period refs, backdated 7 years
    doc = Document()
    doc.add_heading("GDPR Compliance Policy", 0)
    doc.add_paragraph("Review Date: 2019-05-25\nEffective Date: 2018-05-25")
    _h(doc, "Purpose")
    doc.add_paragraph(
        "Ensures compliance with the GDPR transition period requirements and supersedes "
        "obligations under the Data Protection Directive 95/46/EC. Staff must complete "
        "training before the GDPR transition period ends on 25 May 2018."
    )
    _h(doc, "Scope"); doc.add_paragraph("All EU operations and data processors.")
    _h(doc, "Effective Date"); doc.add_paragraph("25 May 2018")
    _h(doc, "Review Date"); doc.add_paragraph("25 May 2019")
    doc.save(out / "gdpr_compliance_policy_2018.docx")
    _backdate(out / "gdpr_compliance_policy_2018.docx", 2700)

    # IT AUP — Windows Server 2012, no owner, backdated 6 years
    doc = Document()
    doc.add_heading("IT Acceptable Use Policy", 0)
    _h(doc, "Purpose")
    doc.add_paragraph(
        "Governs acceptable use of company IT resources. Systems are managed through "
        "Active Directory on Windows Server 2012. Remote access via VPN per Server 2008 "
        "group policy templates."
    )
    _h(doc, "Scope"); doc.add_paragraph("All company-owned and personal devices.")
    _h(doc, "Effective Date"); doc.add_paragraph("1 March 2018")
    _h(doc, "Review Date"); doc.add_paragraph("1 March 2019")
    doc.save(out / "it_acceptable_use_policy_2018.docx")
    _backdate(out / "it_acceptable_use_policy_2018.docx", 2500)

    # Legacy NDA — missing counterparty and expiry, backdated 5 years
    doc = Document()
    doc.add_heading("Non-Disclosure Agreement — Vendor Template", 0)
    _h(doc, "Parties")
    doc.add_paragraph("Between Acme Corp ('Disclosing Party') and [VENDOR NAME TBD] ('Receiving Party').")
    _h(doc, "Term"); doc.add_paragraph("This agreement is effective from 2019-11-01.")
    _h(doc, "Effective Date"); doc.add_paragraph("1 November 2019")
    _h(doc, "Termination"); doc.add_paragraph("30 days written notice by either party.")
    _h(doc, "Signatures"); doc.add_paragraph("Signed counterparts accepted.")
    # Deliberately missing: Owner, Classification, counterparty governance, expiry
    doc.save(out / "vendor_nda_template_2019.docx")
    _backdate(out / "vendor_nda_template_2019.docx", 2100)

    # Project charter — missing Sponsor, Budget, old
    doc = Document()
    doc.add_heading("Project Charter — Legacy BI Modernisation", 0)
    _h(doc, "Objective")
    doc.add_paragraph("Migrate legacy Cognos reports to a modern BI platform.")
    _h(doc, "Scope")
    doc.add_paragraph("Finance and Operations reporting. Excludes HR.")
    _h(doc, "Milestones")
    doc.add_paragraph("• 2020-Q3: Vendor selection\n• 2021-Q1: Migration complete")
    # Missing: Sponsor, Budget — required sections for project_charter
    doc.save(out / "project_charter_legacy_bi_2020.docx")
    _backdate(out / "project_charter_legacy_bi_2020.docx", 1800)

    # Stale status report — 5 years old (threshold 30 days)
    doc = Document()
    doc.add_heading("Weekly Status Report — Week of 9 March 2020", 0)
    _h(doc, "Period"); doc.add_paragraph("9 March – 13 March 2020")
    _h(doc, "Status"); doc.add_paragraph("AMBER — COVID-19 impact on office attendance.")
    _h(doc, "Risks"); doc.add_paragraph("Remote work tooling not yet deployed for all staff.")
    _h(doc, "Next Steps"); doc.add_paragraph("Finalise remote access rollout by 20 March 2020.")
    doc.save(out / "status_report_2020_w11.docx")
    _backdate(out / "status_report_2020_w11.docx", 1900)

    # Stale OKR — backdated 1000 days (threshold 90 days)
    _okr(out / "okr_q3_2022.docx",
         "OKRs Q3 2022 — Data & Analytics",
         "CDO", "internal",
         "Deliver the first version of the centralised data warehouse.",
         ["Launch data warehouse MVP by 30 September 2022.",
          "Onboard Finance reporting by Q3 end.",
          "Reduce manual reporting effort by 40%."],
         "Q3 2022 (July – September)")
    _backdate(out / "okr_q3_2022.docx", 1000)

    # Stale PCI compliance xlsx — retired standard refs, no governance, backdated
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "PCI Compliance Checklist"
    wb.properties.title = "PCI DSS 3.2.1 Compliance Checklist"
    # No creator set
    ws.cell(1, 1, "PCI DSS 3.2.1 Compliance Checklist — FY2020").font = Font(bold=True)
    ws.cell(2, 1, "Note: assessed against PCI-DSS v3 requirements. PCI DSS 3.2.1 controls applied.")
    for col, h in enumerate(["Requirement", "Status", "Evidence", "Owner"], 1):
        ws.cell(4, col, h).font = Font(bold=True)
    reqs = [("Req 1: Firewall", "Compliant", "FW config review Q1 2020", ""),
            ("Req 2: Defaults",  "Non-compliant", "Legacy device found", ""),
            ("Req 6: Secure dev","Compliant", "Code review logs", ""),
            ("Req 8: Access",    "Compliant", "IAM audit 2020-02", "")]
    for i, row in enumerate(reqs, 5):
        for j, v in enumerate(row, 1):
            ws.cell(i, j, v)
    wb.save(out / "pci_compliance_checklist_2020.xlsx")
    _backdate(out / "pci_compliance_checklist_2020.xlsx", 1800)

    # Stale AI ethics pptx — "proposed AI Act", backdated 3 years
    prs = Presentation()
    prs.core_properties.title = "AI Ethics & Regulation Briefing"
    # No author set (no owner)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "AI Ethics & Regulation Briefing"
    slide.placeholders[1].text = "Prepared: October 2021"
    layout2 = prs.slide_layouts[1]
    for title_text, bullets in [
        ("Regulatory Landscape",
         ["Proposed AI Act (2021 draft) — risk-based approach",
          "Draft AI Regulation requires high-risk system documentation",
          "AI Act proposal 2021: Article 13 transparency requirements",
          "GDPR intersection with automated decision-making"]),
        ("Our Approach",
         ["Align with draft AI Regulation classifications",
          "High-risk AI use cases inventoried",
          "Ethics board to review by Q1 2022 — TODO"]),
    ]:
        sl = prs.slides.add_slide(layout2)
        sl.shapes.title.text = title_text
        tf = sl.placeholders[1].text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            tf.add_paragraph().text = b
    prs.save(out / "ai_regulation_briefing_2021.pptx")
    _backdate(out / "ai_regulation_briefing_2021.pptx", 1300)

    # ────────────────────────────────────────────────────────────────────────
    # FLAGGED — GOVERNANCE FAILURES / DRAFT (3)
    # ────────────────────────────────────────────────────────────────────────

    # Employee data dictionary — no owner, no steward, no classification
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Employee Data Dictionary"
    wb.properties.title = "Employee Data Dictionary Draft"
    # No creator, no governance rows
    for col, h in enumerate(["Field Name", "Data Type", "Description"], 1):
        ws.cell(1, col, h).font = Font(bold=True)
    emp_rows = [
        ("employee_id",  "UUID",        "Unique employee identifier."),
        ("full_name",    "VARCHAR(255)","Full legal name."),
        ("salary",       "DECIMAL",     "Annual gross salary EUR."),
        ("performance",  "VARCHAR(10)", "Last review rating."),
        ("manager_id",   "UUID",        "FK to manager employee_id."),
    ]
    for i, row in enumerate(emp_rows, 2):
        for j, v in enumerate(row, 1):
            ws.cell(i, j, v)
    wb.save(out / "employee_data_dictionary_draft.xlsx")

    # Onboarding procedure — placeholders, DRAFT, no owner
    doc = Document()
    doc.add_heading("Employee Onboarding Procedure", 0)
    doc.add_paragraph("DRAFT — DO NOT DISTRIBUTE")
    _h(doc, "Purpose"); doc.add_paragraph("TBD — to be finalised after HR review.")
    _h(doc, "Scope"); doc.add_paragraph("[INSERT scope — all new starters?]")
    _h(doc, "Steps")
    doc.add_paragraph(
        "1. Send welcome email — template at [INSERT link].\n"
        "2. Provision accounts in [INSERT system name]. TODO.\n"
        "3. Assign mandatory training — list TBD.\n"
        "4. 30/60/90-day check-ins — format TODO."
    )
    # Deliberately omit Owner heading and Last Updated
    doc.save(out / "onboarding_procedure_draft.docx")

    # ────────────────────────────────────────────────────────────────────────
    # FLAGGED — ORPHAN FORMATS (.txt / .json / .csv) — 4 files
    # ────────────────────────────────────────────────────────────────────────

    _txt(out / "fraud_risk_assessment.txt",
         "Fraud Risk Assessment — Q1 2025\n\n"
         "Scope: Card-not-present transactions above €500.\n"
         "Current fraud loss rate: 0.04% of revenue.\n"
         "Top risk: account takeover via credential stuffing.\n"
         "Mitigation: deploy velocity rules in fraud engine.\n"
         "Next review: Q2 2025.\n")

    _txt(out / "api_data_contract_draft.txt",
         "API Data Contract — Payments Service\n\n"
         "Schema: TBD\n"
         "Owner: [INSERT team]\n"
         "SLA: [INSERT SLA]\n"
         "Classification: TBD\n"
         "Retention: TBD\n\n"
         "DRAFT — DO NOT DISTRIBUTE\n"
         "Note: This should be migrated to the standard data contract template.\n")

    _json_file(out / "asana_project_dump_q1_2025.json", [
        {"task_id": "11223344", "name": "Update data retention schedule",
         "assignee": "alice@company.com", "due_date": "2025-03-31",
         "status": "in_progress", "last_modified": "2025-04-02"},
        {"task_id": "55667788", "name": "Migrate legacy BI reports",
         "assignee": None, "due_date": None,
         "status": "open", "last_modified": "2023-01-10"},
        {"task_id": "99001122", "name": "Review ISO 27001:2013 audit findings",
         "assignee": "bob@company.com", "due_date": "2022-06-01",
         "status": "open", "last_modified": "2022-05-01"},
    ])

    _csv_file(out / "bc_vendor_master_extract.csv",
              ["VendorID", "Name", "Country", "PaymentTerms", "LastValidated", "Status"],
              [["V001", "Acme Supplies Ltd",    "GB", "Net30", "2021-04-15", "Active"],
               ["V002", "Global Parts GmbH",   "DE", "Net45", "2019-11-30", "Active"],
               ["V003", "FastLogistics SA",    "FR", "Net30", "2020-08-01", "Inactive"],
               ["V004", "Offshore Solutions",  "SG", "Net60", "",           "Under review"],
               ["V005", "Local Services Co",   "IE", "Net30", "2022-03-10", "Active"]])

    total = len([p for p in out.iterdir() if p.is_file()])
    flagged_count = total - clean_count
    print(f"  [FLAGGED] {flagged_count} files (stale/governance/orphan)")
    print(f"\n  Total: {total} files written to {out.resolve()}")
    print(f"  Clean: {clean_count}  |  Flagged (expected): {flagged_count}"
          f"  |  Target rate: {flagged_count/total*100:.0f}%")


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate demo corpus for AI-Readiness Auditor")
    parser.add_argument("--output", default="domain/demo_corpus/files",
                        help="Output directory (default: domain/demo_corpus/files)")
    args = parser.parse_args()
    generate(args.output)


if __name__ == "__main__":
    main()
