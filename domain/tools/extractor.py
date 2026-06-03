"""
Document extractor — pulls body text + embedded metadata from a file.
Imports for each format are deferred so missing libraries only fail on use.
"""
from __future__ import annotations

import csv
import io
import json as _json
from dataclasses import dataclass, field, asdict
from pathlib import Path


@dataclass
class ExtractedDoc:
    text: str
    embedded_metadata: dict = field(default_factory=dict)
    page_or_item_count: int = 0
    extraction_ok: bool = True
    error: str | None = None


def extract_document(path: str, max_chars: int = 12000, timeout_s: int = 10) -> dict:
    """Extract text + embedded metadata from a single document.

    Runs the format-specific extraction in a worker thread capped at timeout_s
    seconds.  This prevents corrupt, encrypted, or pathologically large files
    from hanging the caller (and, by extension, the Gradio event loop).
    """
    import concurrent.futures as _cf

    with _cf.ThreadPoolExecutor(max_workers=1) as pool:
        fut = pool.submit(_extract_dispatch, path)
        try:
            result = fut.result(timeout=timeout_s)
        except _cf.TimeoutError:
            result = ExtractedDoc(
                text="",
                extraction_ok=False,
                error=f"Extraction timed out after {timeout_s}s "
                      f"(file may be corrupt, encrypted, or too large)",
            )
        except Exception as exc:
            result = ExtractedDoc(text="", extraction_ok=False, error=str(exc))

    result.text = result.text[:max_chars]
    return asdict(result)


def _extract_dispatch(path: str) -> ExtractedDoc:
    """Dispatch to the correct format-specific extractor (runs inside a thread)."""
    ext = Path(path).suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext in (".docx", ".doc"):
            return _extract_docx(path)
        if ext in (".xlsx", ".xls"):
            return _extract_xlsx(path)
        if ext in (".pptx", ".ppt"):
            return _extract_pptx(path)
        if ext == ".csv":
            return _extract_csv(path)
        if ext == ".json":
            return _extract_json(path)
        if ext in (".txt", ".md"):
            return _extract_text(path)
        return ExtractedDoc(
            text="",
            extraction_ok=False,
            error=f"Unsupported extension: {ext}",
        )
    except Exception as exc:
        return ExtractedDoc(text="", extraction_ok=False, error=str(exc))


# ── Per-format helpers ────────────────────────────────────────────────────────

def _iso_or_str(value) -> str | None:
    """Coerce datetime / string / None to a string."""
    if value is None:
        return None
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


def _extract_pdf(path: str) -> ExtractedDoc:
    from pypdf import PdfReader
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    text = "\n".join(pages)
    raw_meta = reader.metadata or {}
    meta = {k.lstrip("/"): str(v) for k, v in raw_meta.items() if v}
    return ExtractedDoc(
        text=text,
        embedded_metadata=meta,
        page_or_item_count=len(reader.pages),
    )


def _extract_docx(path: str) -> ExtractedDoc:
    from docx import Document
    doc = Document(path)

    paragraphs = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    table_lines = []
    for table in doc.tables:
        for row in table.rows:
            table_lines.append("\t".join(cell.text for cell in row.cells))
    table_text = "\n".join(table_lines)

    text = paragraphs
    if table_text:
        text += "\n\n[Tables]\n" + table_text

    cp = doc.core_properties
    meta = {
        "author":           cp.author,
        "title":            cp.title,
        "last_modified_by": cp.last_modified_by,
        "created":          _iso_or_str(cp.created),
        "modified":         _iso_or_str(cp.modified),
        "revision":         str(cp.revision) if cp.revision else None,
    }
    meta = {k: v for k, v in meta.items() if v}

    return ExtractedDoc(
        text=text,
        embedded_metadata=meta,
        page_or_item_count=len(doc.paragraphs),
    )


def _extract_xlsx(path: str) -> ExtractedDoc:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    lines = []
    total_rows = 0
    for ws in wb.worksheets:
        lines.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                lines.append("\t".join("" if c is None else str(c) for c in row))
                total_rows += 1
    text = "\n".join(lines)

    props = wb.properties
    meta = {
        "creator":        props.creator,
        "lastModifiedBy": props.lastModifiedBy,
        "created":        _iso_or_str(props.created),
        "modified":       _iso_or_str(props.modified),
        "title":          props.title,
    }
    meta = {k: v for k, v in meta.items() if v}

    return ExtractedDoc(
        text=text,
        embedded_metadata=meta,
        page_or_item_count=total_rows,
    )


def _extract_pptx(path: str) -> ExtractedDoc:
    from pptx import Presentation
    prs = Presentation(path)
    slide_texts = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        if parts:
            slide_texts.append(f"[Slide {i}]\n" + "\n".join(parts))
    text = "\n\n".join(slide_texts)

    cp = prs.core_properties
    meta = {
        "author":           cp.author,
        "last_modified_by": cp.last_modified_by,
        "created":          _iso_or_str(cp.created),
        "modified":         _iso_or_str(cp.modified),
        "title":            cp.title,
    }
    meta = {k: v for k, v in meta.items() if v}

    return ExtractedDoc(
        text=text,
        embedded_metadata=meta,
        page_or_item_count=len(prs.slides),
    )


def _extract_csv(path: str) -> ExtractedDoc:
    with open(path, encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return ExtractedDoc(text="Empty CSV file.", page_or_item_count=0)

    columns = rows[0]
    data_rows = rows[1:]
    sample = data_rows[:5]

    lines = [
        f"CSV file: {len(data_rows)} data rows, {len(columns)} columns.",
        f"Columns: {', '.join(columns)}",
        "Sample rows:",
    ]
    for row in sample:
        lines.append("\t".join(row))

    return ExtractedDoc(
        text="\n".join(lines),
        embedded_metadata={"columns": columns, "row_count": len(data_rows)},
        page_or_item_count=len(data_rows),
    )


def _extract_json(path: str) -> ExtractedDoc:
    with open(path, encoding="utf-8", errors="replace") as f:
        data = _json.load(f)

    if isinstance(data, list):
        count = len(data)
        sample = data[:3]
        text = (
            f"JSON array with {count} items.\n"
            f"Sample (first 3):\n{_json.dumps(sample, indent=2, default=str)}"
        )
    else:
        keys = list(data.keys()) if isinstance(data, dict) else []
        text = (
            f"JSON object with {len(keys)} keys: {', '.join(str(k) for k in keys[:20])}.\n"
            f"{_json.dumps(data, indent=2, default=str)}"
        )
        count = 1

    return ExtractedDoc(
        text=text,
        embedded_metadata={},
        page_or_item_count=count,
    )


def _extract_text(path: str) -> ExtractedDoc:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return ExtractedDoc(
        text=text,
        embedded_metadata={},
        page_or_item_count=text.count("\n") + 1,
    )
